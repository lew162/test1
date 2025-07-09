import os

# 依赖库：pip install pillow pytesseract pdfplumber python-docx openpyxl speechrecognition pydub

try:
    from PIL import Image
    import pytesseract # OCR库
except ImportError:
    Image = None
    pytesseract = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import speech_recognition as sr
    from pydub import AudioSegment
except ImportError:
    sr = None
    AudioSegment = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def file_to_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    if ext in ['.txt', '.md', '.csv', '.json', '.xml']:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff'] and Image and pytesseract:
        image = Image.open(filepath)
        return pytesseract.image_to_string(image)
    elif ext == '.pdf' and pdfplumber:
        text = ''
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ''
        return text
    elif ext in ['.docx'] and docx:
        doc = docx.Document(filepath)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif ext in ['.xlsx', '.xls'] and openpyxl:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        text = ''
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                text += '\t'.join([str(cell) if cell is not None else '' for cell in row]) + '\n'
        return text
    elif ext in ['.wav', '.mp3', '.flac'] and sr and AudioSegment:
        recognizer = sr.Recognizer()
        audio = AudioSegment.from_file(filepath)
        wav_path = filepath + '.temp.wav'
        audio.export(wav_path, format='wav')
        with sr.AudioFile(wav_path) as source:
            audio_data = recognizer.record(source)
            try:
                text = recognizer.recognize_google(audio_data, language='zh-CN')
            except Exception:
                text = ''
        os.remove(wav_path)
        return text
    elif ext in ['.pptx'] and Presentation:
        prs = Presentation(filepath)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    else:
        raise ValueError(f'不支持的文件类型: {ext}')
    
text = file_to_text('D:/example.pdf')  # 替换为你的文件路径
print(text)